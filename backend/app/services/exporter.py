"""
Export Module
Exports valuation results to PowerPoint and PDF formats
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from typing import Dict, Optional
import os
from datetime import datetime


class Exporter:
    """Exports valuation results to various formats"""
    
    def __init__(self, output_dir: str = "outputs"):
        """
        Initialize Exporter
        
        Args:
            output_dir: Directory to save exported files
        """
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
    
    def export_to_powerpoint(
        self,
        ticker: str,
        dcf_results: Optional[Dict] = None,
        comps_results: Optional[Dict] = None,
        wacc_results: Optional[Dict] = None,
        filename: Optional[str] = None
    ) -> str:
        """
        Export valuation results to PowerPoint
        
        Args:
            ticker: Stock ticker symbol
            dcf_results: DCF valuation results
            comps_results: Trading comps results
            wacc_results: WACC calculation results
            filename: Output filename (optional)
            
        Returns:
            Path to created PowerPoint file
        """
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{ticker}_valuation_{timestamp}.pptx"
        
        filepath = os.path.join(self.output_dir, filename)
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"{ticker} Valuation Analysis"
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        
        # WACC slide
        if wacc_results:
            wacc_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
            title_shape = wacc_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
            title_frame = title_shape.text_frame
            title_frame.text = "WACC Calculation"
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            
            # WACC table
            wacc_data = [
                ['Component', 'Value'],
                ['WACC', f"{wacc_results.get('wacc', 0)*100:.2f}%"],
                ['Cost of Equity', f"{wacc_results.get('cost_of_equity', 0)*100:.2f}%"],
                ['Cost of Debt', f"{wacc_results.get('cost_of_debt', 0)*100:.2f}%"],
                ['Equity Weight', f"{wacc_results.get('equity_weight', 0)*100:.2f}%"],
                ['Debt Weight', f"{wacc_results.get('debt_weight', 0)*100:.2f}%"],
                ['Beta', f"{wacc_results.get('beta', 0):.2f}"],
            ]
            
            self._add_table_to_slide(wacc_slide, wacc_data, Inches(1), Inches(1.5), Inches(8), Inches(3))
        
        # DCF slide
        if dcf_results:
            dcf_slide = prs.slides.add_slide(prs.slide_layouts[5])
            title_shape = dcf_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
            title_frame = title_shape.text_frame
            title_frame.text = "DCF Valuation"
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            
            # DCF summary
            dcf_data = [
                ['Metric', 'Value'],
                ['Enterprise Value', f"${dcf_results.get('enterprise_value', 0):,.0f}"],
                ['Equity Value', f"${dcf_results.get('equity_value', 0):,.0f}"],
                ['Fair Value Share Price', f"${dcf_results.get('fair_value_share_price', 0):.2f}"],
                ['WACC', f"{dcf_results.get('wacc', 0)*100:.2f}%"],
                ['Terminal Growth Rate', f"{dcf_results.get('terminal_growth_rate', 0)*100:.2f}%"],
                ['PV of FCFs', f"${dcf_results.get('pv_of_fcfs', 0):,.0f}"],
                ['PV of Terminal Value', f"${dcf_results.get('pv_of_terminal_value', 0):,.0f}"],
            ]
            
            self._add_table_to_slide(dcf_slide, dcf_data, Inches(1), Inches(1.5), Inches(8), Inches(3))
            
            # FCF projections
            if 'projected_fcfs' in dcf_results and not dcf_results['projected_fcfs'].empty:
                fcf_slide = prs.slides.add_slide(prs.slide_layouts[5])
                title_shape = fcf_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
                title_frame = title_shape.text_frame
                title_frame.text = "Projected Free Cash Flows"
                title_frame.paragraphs[0].font.size = Pt(24)
                title_frame.paragraphs[0].font.bold = True
                
                fcf_df = dcf_results['projected_fcfs']
                fcf_data = [['Year', 'FCF ($M)', 'Growth Rate']]
                for _, row in fcf_df.iterrows():
                    fcf_data.append([
                        str(int(row['Year'])),
                        f"${row['FCF']/1e6:.1f}",
                        f"{row['Growth Rate']*100:.1f}%"
                    ])
                
                self._add_table_to_slide(fcf_slide, fcf_data, Inches(1), Inches(1.5), Inches(8), Inches(4))
        
        # Trading Comps slide
        if comps_results:
            comps_slide = prs.slides.add_slide(prs.slide_layouts[5])
            title_shape = comps_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
            title_frame = title_shape.text_frame
            title_frame.text = "Trading Comps Analysis"
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            
            # Valuation ranges
            valuation_ranges = comps_results.get('valuation_ranges', {})
            if valuation_ranges:
                ranges_data = [['Multiple', 'Min Price', 'Median Price', 'Max Price']]
                for multiple_type, ranges in valuation_ranges.items():
                    ranges_data.append([
                        multiple_type,
                        f"${ranges.get('min_share_price', 0):.2f}",
                        f"${ranges.get('median_share_price', 0):.2f}",
                        f"${ranges.get('max_share_price', 0):.2f}"
                    ])
                
                self._add_table_to_slide(comps_slide, ranges_data, Inches(1), Inches(1.5), Inches(8), Inches(3))
        
        prs.save(filepath)
        return filepath
    
    def _add_table_to_slide(self, slide, data, left, top, width, height):
        """Helper method to add a table to a slide"""
        rows = len(data)
        cols = len(data[0]) if data else 0
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        for i, row_data in enumerate(data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                if i == 0:  # Header row
                    cell.text_frame.paragraphs[0].font.bold = True
    
    def export_to_pdf(
        self,
        ticker: str,
        dcf_results: Optional[Dict] = None,
        comps_results: Optional[Dict] = None,
        wacc_results: Optional[Dict] = None,
        filename: Optional[str] = None
    ) -> str:
        """
        Export valuation results to PDF
        
        Args:
            ticker: Stock ticker symbol
            dcf_results: DCF valuation results
            comps_results: Trading comps results
            wacc_results: WACC calculation results
            filename: Output filename (optional)
            
        Returns:
            Path to created PDF file
        """
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{ticker}_valuation_{timestamp}.pdf"
        
        filepath = os.path.join(self.output_dir, filename)
        
        doc = SimpleDocTemplate(filepath, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1f77b4'),
            spaceAfter=30,
            alignment=1  # Center
        )
        story.append(Paragraph(f"{ticker} Valuation Analysis", title_style))
        story.append(Paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # WACC section
        if wacc_results:
            story.append(Paragraph("WACC Calculation", styles['Heading2']))
            wacc_data = [
                ['Component', 'Value'],
                ['WACC', f"{wacc_results.get('wacc', 0)*100:.2f}%"],
                ['Cost of Equity', f"{wacc_results.get('cost_of_equity', 0)*100:.2f}%"],
                ['Cost of Debt', f"{wacc_results.get('cost_of_debt', 0)*100:.2f}%"],
                ['Equity Weight', f"{wacc_results.get('equity_weight', 0)*100:.2f}%"],
                ['Debt Weight', f"{wacc_results.get('debt_weight', 0)*100:.2f}%"],
                ['Beta', f"{wacc_results.get('beta', 0):.2f}"],
            ]
            wacc_table = Table(wacc_data)
            wacc_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            story.append(wacc_table)
            story.append(Spacer(1, 0.3*inch))
        
        # DCF section
        if dcf_results:
            story.append(Paragraph("DCF Valuation", styles['Heading2']))
            dcf_data = [
                ['Metric', 'Value'],
                ['Enterprise Value', f"${dcf_results.get('enterprise_value', 0):,.0f}"],
                ['Equity Value', f"${dcf_results.get('equity_value', 0):,.0f}"],
                ['Fair Value Share Price', f"${dcf_results.get('fair_value_share_price', 0):.2f}"],
                ['WACC', f"{dcf_results.get('wacc', 0)*100:.2f}%"],
                ['Terminal Growth Rate', f"{dcf_results.get('terminal_growth_rate', 0)*100:.2f}%"],
            ]
            dcf_table = Table(dcf_data)
            dcf_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            story.append(dcf_table)
            story.append(Spacer(1, 0.3*inch))
        
        # Trading Comps section
        if comps_results:
            story.append(Paragraph("Trading Comps Analysis", styles['Heading2']))
            valuation_ranges = comps_results.get('valuation_ranges', {})
            if valuation_ranges:
                ranges_data = [['Multiple', 'Min Price', 'Median Price', 'Max Price']]
                for multiple_type, ranges in valuation_ranges.items():
                    ranges_data.append([
                        multiple_type,
                        f"${ranges.get('min_share_price', 0):.2f}",
                        f"${ranges.get('median_share_price', 0):.2f}",
                        f"${ranges.get('max_share_price', 0):.2f}"
                    ])
                
                ranges_table = Table(ranges_data)
                ranges_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(ranges_table)
        
        doc.build(story)
        return filepath

